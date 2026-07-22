# pptx_engine.py
"""
簡報繪製引擎：支援生成「Initial 6 象限」與「Initial vs Re-evaluation 雙期對比 (6+6 頁)」PPTX 簡報。
"""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from typing import Dict, Set, Any
import config

def _rgb(color_tuple):
    return RGBColor(*color_tuple)

def _apply_cell_density(cell):
    cell.margin_top = Inches(config.CELL_INTERNAL_MARGIN)
    cell.margin_bottom = Inches(config.CELL_INTERNAL_MARGIN)
    cell.margin_left = Inches(config.CELL_INTERNAL_MARGIN)
    cell.margin_right = Inches(config.CELL_INTERNAL_MARGIN)

def _set_cell_border_to_white(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    border_xml = (
        f'<a:ln {nsdecls("a")} w="12700" cmpd="s" algn="ctr">'
        f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'<a:prstDash val="solid"/>'
        f'</a:ln>'
    )
    for side in ["lnL", "lnR", "lnT", "lnB"]:
        existing = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{side}")
        if existing is not None:
            tcPr.remove(existing)
        tcPr.append(parse_xml(f'<a:{side} {nsdecls("a")}>{border_xml}</a:{side}>'))

def _set_cell_text(cell, text, size, color, bold=False):
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = config.FONT_PRIMARY
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color

def create_six_sextants_presentation(records: Dict[int, Dict[str, Any]], missing_teeth: Set[int]) -> BytesIO:
    """
    模式 A：生成 6 頁「Initial 單期六象限」標準簡報
    """
    prs = Presentation()
    prs.slide_width = Inches(config.PPT_SLIDE_WIDTH)
    prs.slide_height = Inches(config.PPT_SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]

    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, f"{sextant_name} - Initial Charting", teeth, records, missing_teeth, is_comparison_mode=False)

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream

def create_comparison_presentation(records: Dict[int, Dict[str, Any]], missing_teeth: Set[int]) -> BytesIO:
    """
    模式 B：生成 12 頁「Initial 6 象限 + Initial vs Re-evaluation 對比 6 象限」簡報
    """
    prs = Presentation()
    prs.slide_width = Inches(config.PPT_SLIDE_WIDTH)
    prs.slide_height = Inches(config.PPT_SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]

    # 前 6 頁：Initial 單期
    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, f"{sextant_name} - Initial Stage", teeth, records, missing_teeth, is_comparison_mode=False)

    # 後 6 頁：Initial vs Re-evaluation 對比期
    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, f"{sextant_name} - Initial vs Re-evaluation", teeth, records, missing_teeth, is_comparison_mode=True)

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream

def _draw_sextant_slide(slide, title_text: str, teeth: list, records: Dict[int, Dict[str, Any]], missing_teeth: Set[int], is_comparison_mode: bool):
    """繪製單一象限投影片與表格"""
    # 背景設置純黑
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(config.COLOR_BG_DARK)

    # 標題
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12), Inches(0.8))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = config.FONT_PRIMARY
    p.font.size = Pt(config.FONT_SIZE_TITLE)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = _rgb(config.COLOR_TITLE_WHEAT)

    # 表格設置 (對比模式時，列數會翻倍呈現 I/R)
    row_labels = ["Mobility", "KM", "Probing Depth (B)", "Recession (B)", "Probing Depth (L)", "Recession (L)"]
    rows_per_metric = 2 if is_comparison_mode else 1
    total_rows = 1 + len(row_labels) * rows_per_metric
    total_cols = 1 + len(teeth)

    col_width_left = Inches(2.2)
    col_width_data = Inches(1.8)
    row_height = Inches(config.TABLE_ROW_HEIGHT)

    width = col_width_left + col_width_data * len(teeth)
    height = row_height * total_rows
    left = Inches(0.5)
    top = Inches(1.3)

    table_shape = slide.shapes.add_table(total_rows, total_cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = col_width_left
    for c in range(1, total_cols):
        table.columns[c].width = col_width_data

    # 填入 Header
    _set_cell_text(table.cell(0, 0), "Tooth", config.FONT_SIZE_LABEL, _rgb(config.COLOR_TEXT_WHITE), bold=True)
    for c_i, tooth in enumerate(teeth):
        _set_cell_text(table.cell(0, c_i + 1), str(tooth), config.FONT_SIZE_LABEL, _rgb(config.COLOR_TEXT_WHITE), bold=True)

    # 填入內容與邊框處置
    for r in range(total_rows):
        for c in range(total_cols):
            cell = table.cell(r, c)
            _apply_cell_density(cell)
            _set_cell_border_to_white(cell)
            if r > 0 and c == 0:
                metric_idx = (r - 1) // rows_per_metric
                stage_tag = " (I)" if (is_comparison_mode and (r - 1) % 2 == 0) else (" (R)" if is_comparison_mode else "")
                _set_cell_text(cell, f"{row_labels[metric_idx]}{stage_tag}", 11, _rgb(config.COLOR_TEXT_WHITE), bold=True)
