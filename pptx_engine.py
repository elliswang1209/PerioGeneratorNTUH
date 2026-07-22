# pptx_engine.py

from io import BytesIO
from typing import Dict, Set, Any, List
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

import config
from core_parser import (
    find_tooth_rows, 
    get_tooth_start_columns, 
    collect_comparison_row_indices, 
    get_three_digit_raw_list,
    clean_cell,
    find_furcation_rows
)

VALID_FURCATION_TEETH = {
    18: ["M", "B", "D"], 17: ["M", "B", "D"], 16: ["M", "B", "D"],
    26: ["M", "B", "D"], 27: ["M", "B", "D"], 28: ["M", "B", "D"],
    14: ["M", "D"], 24: ["M", "D"],
    36: ["B", "L"], 37: ["B", "L"], 38: ["B", "L"],
    46: ["B", "L"], 47: ["B", "L"], 48: ["B", "L"]
}

GRADE_MAP = {
    "1": "1", "I": "1",
    "2": "2", "II": "2",
    "3": "3", "III": "3"
}

# 橘黃色設定
COLOR_STAGE_R = RGBColor(255, 180, 0)

def parse_tooth_furcation(df, tooth: int, col: int, f_info: dict) -> str:
    if tooth not in VALID_FURCATION_TEETH or f_info is None:
        return "-"

    target_sites = VALID_FURCATION_TEETH[tooth]
    label_row = f_info["label_row"]
    value_row = f_info["value_row"]

    extracted = {}
    for offset in range(3):
        c = col + offset
        if c < df.shape[1]:
            lbl = clean_cell(df.iloc[label_row, c]).upper()
            val = clean_cell(df.iloc[value_row, c]).upper()
            if lbl in ["M", "B", "D", "L"] and val in GRADE_MAP:
                extracted[lbl] = GRADE_MAP[val]

    result_parts = [f"{s}{extracted[s]}" for s in target_sites if s in extracted]
    return "".join(result_parts) if result_parts else "-"

def parse_tooth_mobility(df, r_idx, col) -> str:
    """跨 3 欄掃描牙齒的 Mobility 數據，精確支援 I, II, III, 1, 2, 3 等格式"""
    if r_idx is None or pd.isna(r_idx) or int(r_idx) >= len(df):
        return "WNL"
    
    mob_map = {
        "1": "I", "I": "I",
        "2": "II", "II": "II",
        "3": "III", "III": "III"
    }
    
    for offset in range(3):
        c = col + offset
        if c < df.shape[1]:
            val = clean_cell(df.iloc[int(r_idx), c]).upper()
            if val in mob_map:
                return mob_map[val]
            elif val in ["WNL", "0", "-"]:
                return "WNL"
            elif val:
                return val
    return "WNL"

def _rgb(color_tuple):
    return RGBColor(*color_tuple)

def _apply_cell_density(cell):
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    cell.margin_left = Inches(0.02)
    cell.margin_right = Inches(0.02)

def _set_cell_border_custom(cell, left_hex="FFFFFF", right_hex="FFFFFF", top_hex="FFFFFF", bottom_hex="FFFFFF", width_pt=1.0):
    tcPr = cell._tc.get_or_add_tcPr()
    w_val = str(int(width_pt * 12700))
    
    sides = {
        'lnL': left_hex,
        'lnR': right_hex,
        'lnT': top_hex,
        'lnB': bottom_hex
    }
    
    for side, hex_color in sides.items():
        existing = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{side}')
        if existing is not None:
            tcPr.remove(existing)
        
        border_xml = (
            f'<a:{side} {nsdecls("a")} w="{w_val}" cmpd="s" algn="ctr">'
            f'  <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'  <a:prstDash val="solid"/>'
            f'</a:{side}>'
        )
        tcPr.append(parse_xml(border_xml))

def _add_diagonal_strikethrough(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    diagonal_xml = (
        f'<a:lnDiagonalDown {nsdecls("a")} w="12700" cmpd="s" algn="ctr">'
        f'  <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
        f'  <a:prstDash val="solid"/>'
        f'</a:lnDiagonalDown>'
    )
    existing = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}lnDiagonalDown')
    if existing is not None:
        tcPr.remove(existing)
    tcPr.append(parse_xml(diagonal_xml))

def _remove_default_table_style(table):
    try:
        tblPr = table._tbl.tblPr
        tableStyleId = tblPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId')
        if tableStyleId is not None:
            tblPr.remove(tableStyleId)
        tblPr.append(parse_xml(f'<a:tableStyleId {nsdecls("a")}>{{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}}</a:tableStyleId>'))
    except Exception:
        pass

def create_six_sextants_presentation(df, missing_teeth: Set[int]) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(config.PPT_SLIDE_WIDTH)
    prs.slide_height = Inches(config.PPT_SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]

    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, sextant_name, teeth, df, missing_teeth, is_comparison=False)

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream

def create_comparison_presentation(df, missing_teeth: Set[int]) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(config.PPT_SLIDE_WIDTH)
    prs.slide_height = Inches(config.PPT_SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]

    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, sextant_name, teeth, df, missing_teeth, is_comparison=False)

    for sextant_name, teeth in config.SEXTANTS.items():
        slide = prs.slides.add_slide(blank_layout)
        _draw_sextant_slide(slide, sextant_name, teeth, df, missing_teeth, is_comparison=True)

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream

def _draw_sextant_slide(slide, title_text: str, teeth: List[int], df, missing_teeth: Set[int], is_comparison: bool):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(config.COLOR_BG_DARK)

    # 1. 投影片主標題 (例如 Upper Right Sextant)
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12), Inches(0.8))
    p = tx_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.name = config.FONT_PRIMARY
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.italic = True
    p.font.underline = True
    p.font.color.rgb = _rgb(config.COLOR_TITLE_WHEAT)
    p.alignment = PP_ALIGN.LEFT

    tooth_rows = find_tooth_rows(df)
    if len(tooth_rows) < 2: return
    upper_cols = get_tooth_start_columns(df, tooth_rows[0])
    lower_cols = get_tooth_start_columns(df, tooth_rows[-1])

    row_map = collect_comparison_row_indices(df)
    is_upper = (teeth[0] < 30)

    # 2. 項目清單
    if not is_comparison:
        active_row_metas = [
            {"label": "Probing Depth (B)" if is_upper else "Probing Depth (L)", "stage": "I", "r_key": "up_b_pd_i" if is_upper else "lo_l_pd_i", "type": "pd"},
            {"label": "Probing Depth (P)" if is_upper else "Probing Depth (B)", "stage": "I", "r_key": "up_p_pd_i" if is_upper else "lo_b_pd_i", "type": "pd"},
            {"label": "Recession", "stage": "I", "r_key": "up_b_gm_i" if is_upper else "lo_l_gm_i", "type": "gm"},
            {"label": "", "stage": "I", "r_key": "up_p_gm_i" if is_upper else "lo_b_gm_i", "type": "gm"},
            {"label": "Masticatory Mucosa", "stage": "I", "r_key": "up_km_i" if is_upper else "lo_km_i", "type": "km"},
            {"label": "Furcation Involvement", "stage": "", "type": "furc"},
            {"label": "Mobility", "stage": "I", "r_key": "up_mob_i" if is_upper else "lo_mob_i", "type": "mob"},
            {"label": "Prognosis", "stage": "", "type": "prog"}
        ]
    else:
        if is_upper:
            active_row_metas = [
                {"label": "Mobility", "stage": "I", "r_key": "up_mob_i", "type": "mob"},
                {"label": "Mobility", "stage": "R", "r_key": "up_mob_r", "type": "mob"},
                {"label": "Masticatory Mucosa", "stage": "I", "r_key": "up_km_i", "type": "km"},
                {"label": "Masticatory Mucosa", "stage": "R", "r_key": "up_km_r", "type": "km"},
                {"label": "Probing Depth (B)", "stage": "I", "r_key": "up_b_pd_i", "type": "pd"},
                {"label": "Probing Depth (B)", "stage": "R", "r_key": "up_b_pd_r", "type": "pd"},
                {"label": "Recession (B)", "stage": "I", "r_key": "up_b_gm_i", "type": "gm"},
                {"label": "Recession (B)", "stage": "R", "r_key": "up_b_gm_r", "type": "gm"},
                {"label": "CAL (B)", "stage": "I", "r_key": "up_b_cal_i", "type": "cal"},
                {"label": "CAL (B)", "stage": "R", "r_key": "up_b_cal_r", "type": "cal"},
                {"label": "Furcation Involvement", "stage": "", "type": "furc"},
                {"label": "Probing Depth (P)", "stage": "I", "r_key": "up_p_pd_i", "type": "pd"},
                {"label": "Probing Depth (P)", "stage": "R", "r_key": "up_p_pd_r", "type": "pd"},
                {"label": "Recession (P)", "stage": "I", "r_key": "up_p_gm_i", "type": "gm"},
                {"label": "Recession (P)", "stage": "R", "r_key": "up_p_gm_r", "type": "gm"},
                {"label": "CAL (P)", "stage": "I", "r_key": "up_p_cal_i", "type": "cal"},
                {"label": "CAL (P)", "stage": "R", "r_key": "up_p_cal_r", "type": "cal"},
            ]
        else:
            active_row_metas = [
                {"label": "Probing Depth (L)", "stage": "I", "r_key": "lo_l_pd_i", "type": "pd"},
                {"label": "Probing Depth (L)", "stage": "R", "r_key": "lo_l_pd_r", "type": "pd"},
                {"label": "Recession (L)", "stage": "I", "r_key": "lo_l_gm_i", "type": "gm"},
                {"label": "Recession (L)", "stage": "R", "r_key": "lo_l_gm_r", "type": "gm"},
                {"label": "CAL (L)", "stage": "I", "r_key": "lo_l_cal_i", "type": "cal"},
                {"label": "CAL (L)", "stage": "R", "r_key": "lo_l_cal_r", "type": "cal"},
                {"label": "Mobility", "stage": "I", "r_key": "lo_mob_i", "type": "mob"},
                {"label": "Mobility", "stage": "R", "r_key": "lo_mob_r", "type": "mob"},
                {"label": "Masticatory Mucosa", "stage": "I", "r_key": "lo_km_i", "type": "km"},
                {"label": "Masticatory Mucosa", "stage": "R", "r_key": "lo_km_r", "type": "km"},
                {"label": "Probing Depth (B)", "stage": "I", "r_key": "lo_b_pd_i", "type": "pd"},
                {"label": "Probing Depth (B)", "stage": "R", "r_key": "lo_b_pd_r", "type": "pd"},
                {"label": "Recession (B)", "stage": "I", "r_key": "lo_b_gm_i", "type": "gm"},
                {"label": "Recession (B)", "stage": "R", "r_key": "lo_b_gm_r", "type": "gm"},
                {"label": "CAL (B)", "stage": "I", "r_key": "lo_b_cal_i", "type": "cal"},
                {"label": "CAL (B)", "stage": "R", "r_key": "lo_b_cal_r", "type": "cal"},
                {"label": "Furcation Involvement", "stage": "", "type": "furc"},
            ]

    total_rows = 1 + len(active_row_metas)
    has_stage_col = is_comparison
    total_cols = (2 if has_stage_col else 1) + len(teeth)

    slide_w, slide_h = config.PPT_SLIDE_WIDTH, config.PPT_SLIDE_HEIGHT

    if not is_comparison:
        col_width_label = Inches(1.5)
        col_width_stage = Inches(0)
        col_width_data  = Inches(1.2)
        top_pos = Inches(1.3)
        row_height = Inches(config.TABLE_ROW_HEIGHT)
    else:
        col_width_label = Inches(1.5)
        col_width_stage = Inches(0.5)
        col_width_data  = Inches(0.8)
        top_pos = Inches(0.25)
        bottom_margin = Inches(0.15)
        available_h = Inches(slide_h) - top_pos - bottom_margin
        row_height = available_h / total_rows

    total_table_width = col_width_label + col_width_stage + (col_width_data * len(teeth))
    total_table_height = row_height * total_rows
    left_pos = Inches(slide_w) - total_table_width

    table_shape = slide.shapes.add_table(total_rows, total_cols, left_pos, top_pos, total_table_width, total_table_height)
    table = table_shape.table
    _remove_default_table_style(table)

    table.columns[0].width = int(col_width_label)
    if has_stage_col:
        table.columns[1].width = int(col_width_stage)
        data_start_col = 2
    else:
        data_start_col = 1

    for c in range(data_start_col, total_cols): 
        table.columns[c].width = int(col_width_data)
        
    for r in range(total_rows): 
        table.rows[r].height = int(row_height)

    text_white, text_alert = _rgb(config.COLOR_TEXT_WHITE), _rgb(config.COLOR_TEXT_ALERT)

    # 3. 表頭 (Col 0: Tooth, Col 2~: 牙號)
    c_t0 = table.cell(0, 0); c_t0.vertical_anchor = MSO_ANCHOR.MIDDLE
    _apply_cell_density(c_t0); c_t0.fill.background()
    p_t0 = c_t0.text_frame.paragraphs[0]; p_t0.alignment = PP_ALIGN.CENTER
    r_t0 = p_t0.add_run(); r_t0.text = "Tooth"
    r_t0.font.name = config.FONT_PRIMARY
    r_t0.font.size = Pt(14)
    r_t0.font.bold = True
    r_t0.font.color.rgb = text_white

    if has_stage_col:
        c_st = table.cell(0, 1); c_st.vertical_anchor = MSO_ANCHOR.MIDDLE
        _apply_cell_density(c_st); c_st.fill.background()

    for c_i, t in enumerate(teeth):
        cell_col_idx = data_start_col + c_i
        c_t = table.cell(0, cell_col_idx); c_t.vertical_anchor = MSO_ANCHOR.MIDDLE
        _apply_cell_density(c_t); c_t.fill.background()
        p_t = c_t.text_frame.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
        r_t = p_t.add_run(); r_t.text = str(t)
        r_t.font.name = config.FONT_PRIMARY
        r_t.font.size = Pt(14)
        r_t.font.bold = True
        r_t.font.color.rgb = text_white

    f_rows = find_furcation_rows(df)
    dynamic_missing = set(missing_teeth)

    # 用於記錄各顆牙齒在 Re-eval 階段是否有 PD >= 5
    teeth_with_high_pd_in_reeval = set()

    # 4. 資料列填寫
    for r_i, meta in enumerate(active_row_metas):
        r_ppt = r_i + 1

        # Label (Col 0)
        c_lbl = table.cell(r_ppt, 0); c_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
        _apply_cell_density(c_lbl); c_lbl.fill.background()
        lbl_text = meta["label"]
        if lbl_text:
            p_lbl = c_lbl.text_frame.paragraphs[0]; p_lbl.alignment = PP_ALIGN.CENTER
            r_lbl = p_lbl.add_run(); r_lbl.text = lbl_text
            r_lbl.font.name = config.FONT_PRIMARY
            r_lbl.font.size = Pt(14) if is_comparison else Pt(12)
            r_lbl.font.bold = True
            r_lbl.font.color.rgb = text_white

        # Stage (Col 1)
        if has_stage_col:
            c_stage = table.cell(r_ppt, 1); c_stage.vertical_anchor = MSO_ANCHOR.MIDDLE
            _apply_cell_density(c_stage); c_stage.fill.background()
            stage_text = meta.get("stage", "")
            if stage_text:
                p_stage = c_stage.text_frame.paragraphs[0]; p_stage.alignment = PP_ALIGN.CENTER
                r_stage = p_stage.add_run(); r_stage.text = stage_text
                r_stage.font.name = config.FONT_PRIMARY
                r_stage.font.size = Pt(13)
                r_stage.font.bold = True
                
                if stage_text.upper() == "R":
                    r_stage.font.color.rgb = COLOR_STAGE_R
                else:
                    r_stage.font.color.rgb = text_white

        # Values (Col 2~)
        for c_i, t in enumerate(teeth):
            cell_col_idx = data_start_col + c_i
            cell_d = table.cell(r_ppt, cell_col_idx); cell_d.vertical_anchor = MSO_ANCHOR.MIDDLE
            _apply_cell_density(cell_d); cell_d.fill.background()

            col = upper_cols.get(t) if is_upper else lower_cols.get(t)
            p_d = cell_d.text_frame.paragraphs[0]; p_d.alignment = PP_ALIGN.CENTER

            if col is not None:
                m_type = meta["type"]
                r_idx = row_map.get(meta.get("r_key"))
                stage_text = meta.get("stage", "")

                if m_type in ["pd", "gm", "cal"]:
                    digits = get_three_digit_raw_list(df, r_idx, col)
                    if m_type == "pd" and all(d.strip() in ['?', ''] for d in digits):
                        dynamic_missing.add(t)

                    for d in digits:
                        display_val = clean_cell(d)
                        run = p_d.add_run(); run.text = f" {display_val} " if len(display_val)>=2 else display_val
                        run.font.name = config.FONT_PRIMARY
                        run.font.size = Pt(12) if is_comparison else Pt(14)
                        
                        if m_type == "pd" and display_val.isdigit() and int(display_val) >= 5:
                            run.font.bold = True; run.font.color.rgb = text_alert
                            if is_comparison and stage_text.upper() == "R":
                                teeth_with_high_pd_in_reeval.add(t)
                        else:
                            run.font.bold = False; run.font.color.rgb = text_white

                elif m_type == "km":
                    val = clean_cell(df.iloc[r_idx, col]) if r_idx is not None else "0"
                    run = p_d.add_run(); run.text = val if val!="" else "0"
                    run.font.name = config.FONT_PRIMARY
                    run.font.size = Pt(12) if is_comparison else Pt(14)
                    run.font.bold, run.font.color.rgb = False, text_white

                elif m_type == "furc":
                    f_info = f_rows[0] if (is_upper and len(f_rows) >= 1) else (f_rows[-1] if len(f_rows) >= 2 else None)
                    text = parse_tooth_furcation(df, t, col, f_info)
                    run = p_d.add_run(); run.text = text
                    run.font.name = config.FONT_PRIMARY
                    run.font.size = Pt(12) if is_comparison else Pt(14)
                    run.font.bold, run.font.color.rgb = False, text_white

                elif m_type == "mob":
                    val = parse_tooth_mobility(df, r_idx, col)
                    run = p_d.add_run(); run.text = val
                    run.font.name = config.FONT_PRIMARY
                    run.font.size = Pt(12) if is_comparison else Pt(14)
                    run.font.bold, run.font.color.rgb = False, text_white

                elif m_type == "prog":
                    run = p_d.add_run(); run.text = "G"
                    run.font.name = config.FONT_PRIMARY
                    run.font.size = Pt(12) if is_comparison else Pt(14)
                    run.font.bold, run.font.color.rgb = False, _rgb((0, 255, 0))

    # 5. 缺牙大合併與斜線
    for c_i, t in enumerate(teeth):
        if t in dynamic_missing:
            cell_col_idx = data_start_col + c_i
            start_cell = table.cell(1, cell_col_idx)
            end_cell = table.cell(total_rows - 1, cell_col_idx)
            start_cell.merge(end_cell)
            
            merged = table.cell(1, cell_col_idx)
            merged.vertical_anchor = MSO_ANCHOR.MIDDLE
            merged.text_frame.clear()
            _apply_cell_density(merged)
            merged.fill.background()
            _add_diagonal_strikethrough(merged)

    # 6. 自動合併相鄰且文字相同的第 0 欄標題列 (Vertical Merge)
    if is_comparison:
        r = 1
        while r < total_rows - 1:
            cell_curr = table.cell(r, 0)
            cell_next = table.cell(r + 1, 0)
            
            txt_curr = cell_curr.text_frame.text.strip()
            txt_next = cell_next.text_frame.text.strip()
            
            if txt_curr and txt_curr == txt_next:
                cell_curr.merge(cell_next)
                
                merged_cell = table.cell(r, 0)
                merged_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                merged_cell.text_frame.clear()
                _apply_cell_density(merged_cell)
                merged_cell.fill.background()
                
                p_m = merged_cell.text_frame.paragraphs[0]
                p_m.alignment = PP_ALIGN.CENTER
                r_m = p_m.add_run()
                r_m.text = txt_curr
                r_m.font.name = config.FONT_PRIMARY
                r_m.font.size = Pt(14)
                r_m.font.bold = True
                r_m.font.color.rgb = text_white
                
                r += 2
            else:
                r += 1

    # 7. 邊框刷新 (若在 R 階段仍有 PD >= 5mm，整欄畫正紅色邊框；否則繪製純白邊框)
    RED_HEX = "FF0000"
    WHITE_HEX = "FFFFFF"

    for c_i, t in enumerate(teeth):
        cell_col = data_start_col + c_i
        is_high_pd_tooth = (is_comparison and t in teeth_with_high_pd_in_reeval and t not in dynamic_missing)

        left_color = RED_HEX if is_high_pd_tooth else WHITE_HEX
        right_color = RED_HEX if is_high_pd_tooth else WHITE_HEX
        border_width = 2.0 if is_high_pd_tooth else 1.0

        for r in range(total_rows):
            top_color = RED_HEX if (is_high_pd_tooth and r == 0) else WHITE_HEX
            bottom_color = RED_HEX if (is_high_pd_tooth and r == total_rows - 1) else WHITE_HEX

            _set_cell_border_custom(
                table.cell(r, cell_col),
                left_hex=left_color,
                right_hex=right_color,
                top_hex=top_color,
                bottom_hex=bottom_color,
                width_pt=border_width
            )

    # 第 0 欄與 Stage 欄維持純白邊框
    for r in range(total_rows):
        _set_cell_border_custom(table.cell(r, 0), WHITE_HEX, WHITE_HEX, WHITE_HEX, WHITE_HEX, 1.0)
        if has_stage_col:
            _set_cell_border_custom(table.cell(r, 1), WHITE_HEX, WHITE_HEX, WHITE_HEX, WHITE_HEX, 1.0)
